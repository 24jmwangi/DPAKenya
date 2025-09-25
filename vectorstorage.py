import os
import faiss
import pickle
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from sklearn.preprocessing import normalize
import numpy as np

# Load local embedding model
model = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def chunk_text(text, max_tokens=300):
    words = text.split()
    chunks, chunk = [], []
    for word in words:
        chunk.append(word)
        if len(chunk) >= max_tokens:
            chunks.append(" ".join(chunk))
            chunk = []
    if chunk:
        chunks.append(" ".join(chunk))
    return chunks

def embed_text(chunks):
    embeddings = model.encode(chunks, convert_to_numpy=True, normalize_embeddings=True)
    return embeddings.astype("float32")

def build_faiss_index(files):
    all_chunks, metadata = [], []
    for file in files:
        text = extract_text_from_pptx(file)
        chunks = chunk_text(text)
        all_chunks.extend(chunks)
        metadata.extend([{"source": file, "text": c} for c in chunks])

    embeddings = embed_text(all_chunks)

    dim = embeddings.shape[1]
    index = faiss.IndexFlatIP(dim)  # cosine similarity
    index.add(embeddings)

    return index, metadata

def save_vectorstore(index, metadata, path="vectorstore.pkl"):
    with open(path, "wb") as f:
        pickle.dump({"index": index, "metadata": metadata}, f)

if __name__ == "__main__":
    pptx_files = [
        "globalp.pptx",
        "glossary.pptx",
        "m1legalaspects.pptx",
        "m2borgani.pptx",
        "m2technical.pptx",
        "m3privag.pptx"
    ]
    index, metadata = build_faiss_index(pptx_files)
    save_vectorstore(index, metadata)
    print("Vectorstore built locally")
